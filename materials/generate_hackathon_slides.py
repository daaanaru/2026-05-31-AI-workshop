#!/usr/bin/env python3
"""ハッカソンアイデア25テーマ — 簡易スライド生成（50枚）"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# --- Colors (既存スクリプトと同じパレット) ---
MOSS_DEEP = RGBColor(0x2B, 0x3A, 0x2E)  # ダークグリーン
PAPER = RGBColor(0xF5, 0xF2, 0xEB)     # 和紙色
INK = RGBColor(0x2D, 0x2D, 0x2D)
ACCENT_GREEN = RGBColor(0x8B, 0xA8, 0x6C)
SNOW = RGBColor(0xFA, 0xFA, 0xF5)
ROUTE = RGBColor(0xC1, 0x5A, 0x2B)     # ルート色（オレンジ系）
INK_SOFT = RGBColor(0x4F, 0x62, 0x51)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FADED = RGBColor(0x8C, 0x9E, 0x7A)

# Themes data
themes = [
    {"num": 1, "title": "PRレビュー要約Bot", "subtitle": "PRのdiffをAIが要約、レビューポイントを提示", "difficulty": "★★☆", "tech": "Python/Node.js + GitHub API + LLM API", "problem": "PRが大きくなるとレビューに時間がかかる。要点だけ先に知りたい", "solution": "GitHub PR URL入力→diff取得→LLMが変更サマリー+レビューポイントを生成", "mvp": ["PR URL入力フォーム", "diff取得→LLM要約", "結果をMarkdownで表示"]},
    {"num": 2, "title": "議事録→TODO抽出マシン", "subtitle": "会議テキストから決定事項・TODO・担当者を自動抽出", "difficulty": "★☆☆", "tech": "HTML + JS + LLM API", "problem": "会議後に「誰が何やるんだっけ？」ってなる", "solution": "議事録テキスト入力→LLMが構造化→決定事項/TODO/担当を分離出力", "mvp": ["テキスト入力エリア", "LLM構造化処理", "Markdown/JSON出力"]},
    {"num": 3, "title": "ドキュメントQ&A Bot", "subtitle": "自分のドキュメントに質問できるRAGチャットBot", "difficulty": "★★★", "tech": "Python + LangChain + ChromaDB + Streamlit", "problem": "社内Wiki・READMEが膨大で欲しい情報が見つからない", "solution": "Markdownアップロード→ベクトル化→質問に引用元付きで回答", "mvp": ["ファイルアップロード", "RAGパイプライン構築", "チャットUI"]},
    {"num": 4, "title": "AI日報ジェネレーター", "subtitle": "gitログ+チャットから今日の日報を自動生成", "difficulty": "★☆☆", "tech": "Python/Node.js + Git CLI + LLM API", "problem": "毎日の日報作成がめんどくさい", "solution": "git log取得+チャットログ入力→LLMが日報フォーマットで生成", "mvp": ["git log取得", "テキスト入力", "日報フォーマット出力"]},
    {"num": 5, "title": "コマンドライン翻訳家", "subtitle": "自然言語→正しいターミナルコマンドを返すCLI", "difficulty": "★☆☆", "tech": "Python/Node.js + LLM API", "problem": "コマンドのオプションを毎回ググるのがだるい", "solution": "「○○したい」→LLMがコマンド生成→確認→実行", "mvp": ["CLIインターフェース", "LLMコマンド生成", "確認プロンプト+実行"]},
    {"num": 6, "title": "CLAUDE.mdジェネレーター", "subtitle": "リポジトリからCLAUDE.mdを自動生成", "difficulty": "★★☆", "tech": "Python/Node.js + ファイル操作 + LLM API", "problem": "Claude Codeを使いたいけどCLAUDE.mdを書くのが面倒", "solution": "リポジトリ構造・package.json・README解析→CLAUDE.md自動生成", "mvp": ["リポジトリパス指定", "構造解析", "CLAUDE.md出力"]},
    {"num": 7, "title": "スクショ→Issue変換", "subtitle": "バグのスクショから再現手順付きIssueを自動作成", "difficulty": "★★☆", "tech": "HTML + JS + マルチモーダルLLM API", "problem": "スクショだけ飛んできて再現手順が書かれてない", "solution": "画像アップ→マルチモーダルLLMで解析→Issue形式で出力", "mvp": ["画像D&D", "LLM画像解析", "Issueテンプレ出力"]},
    {"num": 8, "title": "プロンプト共有ライブラリ", "subtitle": "チームで使えるプロンプト集+コピーボタン+評価", "difficulty": "★☆☆", "tech": "HTML + CSS + JS", "problem": "いいプロンプト見つけたけど共有場所がない", "solution": "カテゴリ分け+コピーボタン+評価機能のWebアプリ", "mvp": ["SPA UI", "カテゴリ分類", "LocalStorage保存"]},
    {"num": 9, "title": "git diff読み上げBot", "subtitle": "diffを日本語の変更サマリーに変換してチャットに投稿", "difficulty": "★★☆", "tech": "Python/Node.js + Git CLI + LLM API", "problem": "diffが大きいと読む気が失せる", "solution": "git diff取得→LLMが要約→変更理由推測+レビューポイント生成", "mvp": ["diff取得", "LLM要約", "チャット投稿"]},
    {"num": 10, "title": "エラーメッセージ翻訳家", "subtitle": "エラーログを貼ると原因・対処法を日本語で返す", "difficulty": "★☆☆", "tech": "HTML + JS + LLM API", "problem": "英語のスタックトレースが読めない/読みたくない", "solution": "エラー貼付→LLMが原因分析→対処法ステップバイステップ", "mvp": ["テキスト入力", "LLM解析", "対処法表示"]},
    {"num": 11, "title": "GameForge AI", "subtitle": "テキスト→60秒でプレイ可能なブラウザゲーム生成", "difficulty": "★★☆", "tech": "HTML + JS + LLM API", "problem": "ゲームのアイデアはあるが実装する時間がない", "solution": "テキストでゲーム説明→LLMがHTML+JS+CSS生成→iframe内で即プレイ", "mvp": ["テキスト入力", "コード生成", "iframe実行"]},
    {"num": 12, "title": "Prism", "subtitle": "Slack/Discordの雑談からバグ報告・機能要望を自動抽出", "difficulty": "★★☆", "tech": "HTML + JS + LLM API", "problem": "ノイジーなチャットに重要情報が埋もれる", "solution": "チャットログ入力→LLMがカテゴリ分け→Issue/バックログ形式で出力", "mvp": ["テキスト入力", "分類処理", "構造化出力"]},
    {"num": 13, "title": "ASCII（リポジトリ質問Bot）", "subtitle": "リポジトリに何を聞いても答えるBot（MCP活用）", "difficulty": "★★☆", "tech": "Claude Code SDK + GitHub MCP", "problem": "新入社員の「これどうなってるんですか？」が多すぎる", "solution": "リポジトリ指定→コードベースを理解→質問にコード引用付きで回答", "mvp": ["リポジトリ指定", "コード理解", "Q&A応答"]},
    {"num": 14, "title": "cart-to-kitchen", "subtitle": "買い物かごの写真→食材認識→レシピ提案", "difficulty": "★★☆", "tech": "HTML + JS + マルチモーダルLLM API", "problem": "食材はあるけど何を作ればいいかわからない", "solution": "写真アップ→食材リスト生成→3レシピ提案", "mvp": ["画像アップ", "食材認識", "レシピ出力"]},
    {"num": 15, "title": "仕様書＆テスト自動生成", "subtitle": "GitHub+Notionから要件書・テストケースを自動生成", "difficulty": "★★★", "tech": "Python/Node.js + GitHub API + LLM API", "problem": "QAの手作業コピペが多すぎる", "solution": "PR+Notionページ入力→LLMが要件書+テストケース一覧をMarkdownで出力", "mvp": ["URL入力", "複数API連携", "構造化出力"]},
    {"num": 16, "title": "税務会話AI", "subtitle": "確定申告を会話形式で完了するAI", "difficulty": "★☆☆", "tech": "HTML + JS + LLM API", "problem": "確定申告・経費仕分けが面倒", "solution": "「経費として落とせますか？」→AIが質問→仕分け結果を出力", "mvp": ["チャットUI", "質問応答ループ", "仕分け出力"]},
    {"num": 17, "title": "リアルタイム会議マインドマップ", "subtitle": "会議の流れがリアルタイムでマインドマップに育つ", "difficulty": "★★☆", "tech": "HTML + JS + LLM API + Markmap", "problem": "字幕は流れるだけ。会議の全体像が掴めない", "solution": "テキスト逐次入力→LLMが定期的にマインドマップ再生成→リアルタイム描画", "mvp": ["テキスト入力", "LLMマップ生成", "Markmap描画"]},
    {"num": 18, "title": "\"今なに言った？\"キャッチャー", "subtitle": "自分の名前・決定事項・質問の瞬間だけポップアップ通知", "difficulty": "★☆☆", "tech": "Python/JS + 文字起こしAPI + 通知API", "problem": "画面共有を見てると字幕が追えない。視線が1つしかない", "solution": "字幕テキスト監視→キーワード検出→デスクトップ通知", "mvp": ["テキスト監視", "キーワード検出", "通知表示"]},
    {"num": 19, "title": "会議の\"空気\"テキスト化", "subtitle": "笑い・沈黙・声の変化をメタ情報として字幕に付加", "difficulty": "★★★", "tech": "Python + 音声分析API + LLM", "problem": "字幕だけだと冗談なのかマジなのか分からない", "solution": "音声分析→[笑い][沈黙5秒]等を検出→字幕にインライン挿入", "mvp": ["音声アップロード", "感情分析", "メタ情報付き字幕"]},
    {"num": 20, "title": "音の可視化ダッシュボード", "subtitle": "環境音をリアルタイム分類してアイコン表示", "difficulty": "★★☆", "tech": "JS + Web Audio API + TensorFlow.js", "problem": "オフィスで呼ぶ声・電話・アラームが聞こえない", "solution": "マイクで環境音取得→音声分類→アイコン+方向で表示", "mvp": ["マイク入力", "YAMNet分類", "アイコン表示"]},
    {"num": 21, "title": "非同期スタンドアップBot", "subtitle": "朝会を音声なしで完結。AI要約+依存関係検出", "difficulty": "★★☆", "tech": "Python + Slack API + LLM API", "problem": "朝会は音声で回る。リアルタイムで全員の発言を追うのは大変", "solution": "各自テキスト投稿→LLM要約→依存関係検出レポート", "mvp": ["Bot質問", "回答収集", "要約+依存検出"]},
    {"num": 22, "title": "ペアプロ仲介AI", "subtitle": "テキストベースのペアプロをAIが文脈補完", "difficulty": "★★☆", "tech": "HTML + JS + WebSocket + LLM API", "problem": "ペアプロは声が前提。テキストだとテンポが遅い", "solution": "2人のテキスト入力→AIが文脈補完・整理→意図を構造化表示", "mvp": ["2画面チャット", "LLM文脈補完", "構造化表示"]},
    {"num": 23, "title": "見逃し字幕リプレイ", "subtitle": "リアルタイム字幕を巻き戻せる。重要発言ハイライト", "difficulty": "★☆☆", "tech": "HTML + JS + 文字起こしAPI + LLM API", "problem": "字幕は流れていく。目を離したら終わり", "solution": "字幕をバッファ蓄積→「直近30秒再表示」ボタン→重要発言ハイライト", "mvp": ["リングバッファ", "再表示ボタン", "ハイライト"]},
    {"num": 24, "title": "\"聞こえる前提\"検出器", "subtitle": "仕様書の音声前提を検出し代替案を提案", "difficulty": "★☆☆", "tech": "HTML + JS + LLM API", "problem": "仕様書に「電話で通知」と書いてあっても聴覚障害者には届かない", "solution": "テキスト入力→LLMが聞こえる前提を検出→代替案提案", "mvp": ["テキスト入力", "LLM検出", "代替案表示"]},
    {"num": 25, "title": "動画→構造化ノート変換", "subtitle": "YouTube URL→章立て+�条書きMarkdownに変換", "difficulty": "★★☆", "tech": "Python + YouTube API + LLM API", "problem": "技術動画の字幕品質がバラバラ。文字で読みたい", "solution": "URL入力→字幕取得→LLMが章立て+要約+コードスニペット抽出", "mvp": ["URL入力", "字幕取得", "Markdown出力"]},
]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ====================================================================
# Helper functions (既存スクリプトから流用)
# ====================================================================

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=INK, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Noto Sans JP", line_spacing=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if line_spacing:
        p.line_spacing = line_spacing
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_size=18,
                  color=INK, font_name="Noto Sans JP", spacing=Pt(6),
                  bold_first=False, line_spacing=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        if bold_first and i == 0:
            p.font.bold = True
        p.space_after = spacing
        if line_spacing:
            p.line_spacing = line_spacing
    return txBox


def add_section_label(slide, text, top=Inches(0.6), color=ROUTE):
    add_text_box(slide, Inches(0.8), top, Inches(6), Inches(0.4),
                 text, font_size=11, color=color, bold=True,
                 font_name="JetBrains Mono")


def accent_bar(slide, left, top, width=Inches(1.5), height=Pt(4)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ROUTE
    shape.line.fill.background()


def card(slide, left, top, w, h, bg_color=WHITE, border_color=RGBColor(0xCF, 0xD5, 0xB4)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1)
    return shape


# ====================================================================
# Slide generation loop
# ====================================================================

for theme in themes:
    # ---- Slide 1: Title slide (DARK) ----
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, MOSS_DEEP)

    # Theme number + title (large)
    add_text_box(s, Inches(0.8), Inches(0.6), Inches(8), Inches(0.4),
                 f"THEME {theme['num']:02d}", font_size=12,
                 color=ACCENT_GREEN, font_name="JetBrains Mono", bold=True)

    add_text_box(s, Inches(0.8), Inches(1.3), Inches(11), Inches(1.8),
                 theme['title'], font_size=48,
                 color=SNOW, bold=True)

    # Subtitle
    add_text_box(s, Inches(0.8), Inches(3.3), Inches(11), Inches(0.6),
                 theme['subtitle'], font_size=20,
                 color=RGBColor(0xCF, 0xD5, 0xB4))

    # Difficulty + Tech stack (bottom)
    accent_bar(s, Inches(0.8), Inches(4.5), Inches(1.5))

    add_text_box(s, Inches(0.8), Inches(5.0), Inches(3), Inches(0.4),
                 f"難易度: {theme['difficulty']}", font_size=14,
                 color=ACCENT_GREEN, font_name="JetBrains Mono")

    add_text_box(s, Inches(0.8), Inches(5.5), Inches(11), Inches(1.5),
                 f"Tech: {theme['tech']}", font_size=13,
                 color=FADED,
                 font_name="JetBrains Mono")

    # ---- Slide 2: Detail slide (LIGHT) ----
    s = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(s, PAPER)

    add_section_label(s, f"THEME {theme['num']:02d} · DETAIL")

    # Problem card (left)
    card(s, Inches(0.8), Inches(1.3), Inches(5.8), Inches(3.0))
    add_text_box(s, Inches(1.0), Inches(1.5), Inches(5.4), Inches(0.35),
                 "課題（Problem）", font_size=14,
                 color=ROUTE, bold=True)
    add_text_box(s, Inches(1.0), Inches(1.95), Inches(5.4), Inches(2.2),
                 theme['problem'], font_size=14,
                 color=INK)

    # Solution card (right)
    card(s, Inches(7.0), Inches(1.3), Inches(5.8), Inches(3.0))
    add_text_box(s, Inches(7.2), Inches(1.5), Inches(5.4), Inches(0.35),
                 "解決策（Solution）", font_size=14,
                 color=ROUTE, bold=True)
    add_text_box(s, Inches(7.2), Inches(1.95), Inches(5.4), Inches(2.2),
                 theme['solution'], font_size=14,
                 color=INK)

    # MVP section (bottom)
    add_text_box(s, Inches(0.8), Inches(4.5), Inches(11), Inches(0.35),
                 "MVP（3時間で作る範囲）", font_size=14,
                 color=INK, bold=True)

    mvp_text = "\n".join([f"・ {item}" for item in theme['mvp']])
    add_text_box(s, Inches(0.8), Inches(4.95), Inches(11), Inches(2.0),
                 mvp_text, font_size=13,
                 color=INK)


# ====================================================================
# Save
# ====================================================================

out = os.path.join(os.path.dirname(__file__), "hackathon_ideas_slides.pptx")
prs.save(out)
total_slides = len(themes) * 2
print(f"✓ Saved: {out}")
print(f"  Total: {total_slides} slides ({len(themes)} themes × 2)")
