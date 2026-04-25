#!/usr/bin/env python3
"""A1: 上様プレゼン資料 — AIツール遍歴と展望（30分用）v2 肉付け版"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# --- Colors ---
INK = RGBColor(0x1B, 0x2A, 0x1F)
INK_SOFT = RGBColor(0x4F, 0x62, 0x51)
MOSS = RGBColor(0x2F, 0x5F, 0x3A)
MOSS_DEEP = RGBColor(0x1F, 0x41, 0x28)
ROUTE = RGBColor(0xC1, 0x5A, 0x2B)
PAPER = RGBColor(0xEE, 0xF0, 0xE1)
SNOW = RGBColor(0xF4, 0xF2, 0xE0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BARK = RGBColor(0x6B, 0x5A, 0x3A)
CARD_BG = RGBColor(0xF5, 0xF6, 0xEC)
ACCENT_GREEN = RGBColor(0xA8, 0xC0, 0x8A)
FADED = RGBColor(0x8C, 0x9E, 0x7A)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
TOTAL = 15


def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text_box(slide, left, top, width, height, text, font_size=18,
                 color=INK, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Noto Sans JP", line_spacing=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    if line_spacing:
        p.line_spacing = line_spacing
    return txBox


def add_multiline(slide, left, top, width, height, lines, font_size=18,
                  color=INK, font_name="Noto Sans JP", spacing=Pt(6),
                  bold_first=False, line_spacing=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        # Support bold prefix with **text**
        if line.startswith("**") and "**" in line[2:]:
            end = line.index("**", 2)
            bold_part = line[2:end]
            rest = line[end+2:]
            run1 = p.add_run()
            run1.text = bold_part
            run1.font.size = Pt(font_size)
            run1.font.color.rgb = color
            run1.font.name = font_name
            run1.font.bold = True
            if rest:
                run2 = p.add_run()
                run2.text = rest
                run2.font.size = Pt(font_size)
                run2.font.color.rgb = color
                run2.font.name = font_name
                run2.font.bold = False
        else:
            p.text = line
            p.font.size = Pt(font_size)
            p.font.color.rgb = color
            p.font.name = font_name
            if bold_first and i == 0:
                p.font.bold = True
        p.space_after = spacing
        if line_spacing:
            p.line_spacing = line_spacing
    return txBox


def add_section_label(slide, text, top=Inches(0.6), color=ROUTE):
    add_text_box(slide, Inches(0.8), top, Inches(6), Inches(0.4),
                 text, font_size=11, color=color, bold=True,
                 font_name="JetBrains Mono")


def add_slide_number(slide, num):
    add_text_box(slide, Inches(12), Inches(7.0), Inches(1), Inches(0.3),
                 f"{num} / {TOTAL}", font_size=9, color=INK_SOFT,
                 alignment=PP_ALIGN.RIGHT, font_name="JetBrains Mono")


def accent_bar(slide, left, top, width=Inches(1.5), height=Pt(4)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ROUTE
    shape.line.fill.background()


def card(slide, left, top, w, h, bg_color=WHITE, border_color=RGBColor(0xCF, 0xD5, 0xB4)):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    shape.line.color.rgb = border_color
    shape.line.width = Pt(1)
    return shape


# ====================================================================
# SLIDE 1: Title
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, MOSS_DEEP)

add_text_box(s, Inches(0.8), Inches(0.5), Inches(8), Inches(0.4),
             "DEAF ENGINEERS AI WORKSHOP · VOL.01", font_size=12,
             color=ACCENT_GREEN, font_name="JetBrains Mono")

add_text_box(s, Inches(0.8), Inches(2.0), Inches(10), Inches(2.0),
             "AIツール遍歴と\nこれからの使い方", font_size=54,
             color=SNOW, bold=True)

accent_bar(s, Inches(0.8), Inches(4.5), Inches(2))

add_text_box(s, Inches(0.8), Inches(5.0), Inches(6), Inches(0.6),
             "下田 成大（しもだ なるひろ）", font_size=20,
             color=RGBColor(0xCF, 0xD5, 0xB4))

add_text_box(s, Inches(0.8), Inches(5.6), Inches(6), Inches(0.6),
             "2026.05.31 · Deaf Engineers AI Workshop", font_size=14,
             color=FADED, font_name="JetBrains Mono")

# ====================================================================
# SLIDE 2: 自己紹介 (1min)
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, PAPER)
add_section_label(s, "INTRODUCTION")
add_slide_number(s, 2)
add_text_box(s, Inches(11.5), Inches(6.8), Inches(1.5), Inches(0.3),
             "~ 1 min", font_size=9, color=INK_SOFT,
             alignment=PP_ALIGN.RIGHT, font_name="JetBrains Mono")

add_text_box(s, Inches(0.8), Inches(1.2), Inches(10), Inches(0.8),
             "自己紹介", font_size=36, color=INK, bold=True)
accent_bar(s, Inches(0.8), Inches(2.1))

add_multiline(s, Inches(0.8), Inches(2.6), Inches(5.5), Inches(4.5), [
    "下田 成大（だーなる）",
    "",
    "DaNARU 代表（2026年3月設立）",
    "プログラミング初心者",
    "",
    "AIエージェント12体と会社を運営中",
    "社員は全員AI。人間は僕だけ。",
], font_size=20, spacing=Pt(4))

# Right side: CLAUDE.md code block
card(s, Inches(7.2), Inches(2.6), Inches(5.3), Inches(4),
     bg_color=RGBColor(0x1B, 0x2A, 0x1F),
     border_color=RGBColor(0x4F, 0x62, 0x51))
add_text_box(s, Inches(7.4), Inches(2.7), Inches(4.8), Inches(0.3),
             "CLAUDE.md", font_size=10, color=ACCENT_GREEN,
             font_name="JetBrains Mono")
add_multiline(s, Inches(7.4), Inches(3.1), Inches(4.9), Inches(3.3), [
    "# DaNARU \u2014 \u4f1a\u793e\u61b2\u6cd5",
    "",
    "\u793e\u540d: DaNARU\uff08\u3060\u30fc\u306a\u308b\uff09",
    "\u793e\u54e1: \u4ee3\u8868 + AI\u30a8\u30fc\u30b8\u30a7\u30f3\u30c8\u7fa4",
    "\u4e8b\u696d: AI\u3092\u6d3b\u7528\u3057\u305f\u8907\u5408\u53ce\u76ca\u4e8b\u696d",
    "",
    "\u4ee3\u8868\u306e\u4ed5\u4e8b\u306f\u610f\u601d\u6c7a\u5b9a\uff0b\u8a8d\u53ef\u3060\u3051\u3002",
    "\u5b9f\u884c\u306f\u5168\u90e8\u5e55\u5e9c\u304c\u3084\u308b\u3002",
], font_size=13, color=ACCENT_GREEN, font_name="JetBrains Mono", spacing=Pt(2))

# ====================================================================
# SLIDE 3: 今日伝えたいこと (1min) — NEW
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, MOSS_DEEP)
add_slide_number(s, 3)
add_text_box(s, Inches(11.5), Inches(6.8), Inches(1.5), Inches(0.3),
             "~ 1 min", font_size=9, color=FADED,
             alignment=PP_ALIGN.RIGHT, font_name="JetBrains Mono")

add_text_box(s, Inches(0.8), Inches(0.5), Inches(8), Inches(0.4),
             "KEY MESSAGE", font_size=11,
             color=ACCENT_GREEN, bold=True, font_name="JetBrains Mono")

add_text_box(s, Inches(0.8), Inches(1.3), Inches(10), Inches(0.8),
             "今日伝えたいこと", font_size=40, color=SNOW, bold=True)
accent_bar(s, Inches(0.8), Inches(2.3), Inches(2))

add_multiline(s, Inches(0.8), Inches(3.0), Inches(11), Inches(4), [
    "**\u2460 \u30d7\u30ed\u30b0\u30e9\u30df\u30f3\u30b0\u521d\u5fc3\u8005\u3067\u3082\u3001AI\u3067\u4f1a\u793e\u3092\u4f5c\u308c\u305f**",
    "",
    "**\u2461 \u5931\u6557\u3057\u307e\u304f\u3063\u305f\u3051\u3069\u3001\u305d\u308c\u304c\u5168\u90e8\u5b66\u3073\u306b\u306a\u3063\u305f**",
    "",
    "**\u2462 \u300c\u4f55\u3092\u4f5c\u308b\u304b\u300d\u3088\u308a\u300c\u3069\u306e\u554f\u984c\u3092\u89e3\u304f\u304b\u300d**",
], font_size=28, color=SNOW, spacing=Pt(8))

# ====================================================================
# SLIDE 4: AI以前 (2min)
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, PAPER)
add_section_label(s, "CHAPTER 01 · BEFORE AI")
add_slide_number(s, 4)
add_text_box(s, Inches(11.5), Inches(6.8), Inches(1.5), Inches(0.3),
             "~ 2 min", font_size=9, color=INK_SOFT,
             alignment=PP_ALIGN.RIGHT, font_name="JetBrains Mono")

add_text_box(s, Inches(0.8), Inches(1.2), Inches(10), Inches(0.8),
             "AI以前の世界", font_size=36, color=INK, bold=True)
accent_bar(s, Inches(0.8), Inches(2.1))

add_multiline(s, Inches(0.8), Inches(2.6), Inches(11), Inches(4.5), [
    "プログラミング初心者だった。",
    "コードは読めても、自分では書けなかった。",
    "",
    "・ドキュメントを読むのに時間がかかる",
    "・エラーが出ても何が起きたか分からない",
    "・「動くもの」を作るまでが果てしなく遠い",
    "",
    "プログラミングを覚えたかったけど、",
    "独学で壁にぶつかっていた。",
], font_size=20, spacing=Pt(4))

# ====================================================================
# SLIDE 4: 最初の出会い — ChatGPT
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, PAPER)
add_section_label(s, "CHAPTER 02 · FIRST CONTACT — 2022")
add_slide_number(s, 5)
add_text_box(s, Inches(11.5), Inches(6.8), Inches(1.5), Inches(0.3),
             "~ 3 min", font_size=9, color=INK_SOFT,
             alignment=PP_ALIGN.RIGHT, font_name="JetBrains Mono")

add_text_box(s, Inches(0.8), Inches(1.2), Inches(10), Inches(0.8),
             "ChatGPTとの出会い", font_size=36, color=INK, bold=True)
accent_bar(s, Inches(0.8), Inches(2.1))

add_multiline(s, Inches(0.8), Inches(2.6), Inches(11), Inches(4.5), [
    "**2022年11月** — ChatGPT 3.5がリリースされた",
    "",
    "新しもの好きだから飛びついた。",
    "「これまでとは違うチャットボット」に嬉々として使い倒した。",
    "",
    "……でも、いつしか飽きて使わなくなった。",
    "",
    "**2025年春** — GPT-4oの「ジブリ風」が流行った時、また飛びついた。",
    "今度は飽きなかった。毎日使うようになった。",
    "",
    "「キャズム」を超えた瞬間だった。",
], font_size=20, spacing=Pt(4))

# ====================================================================
# SLIDE 5: ツール遍歴タイムライン
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, PAPER)
add_section_label(s, "CHAPTER 03 · TIMELINE")
add_slide_number(s, 6)
add_text_box(s, Inches(11.5), Inches(6.8), Inches(1.5), Inches(0.3),
             "~ 2 min", font_size=9, color=INK_SOFT,
             alignment=PP_ALIGN.RIGHT, font_name="JetBrains Mono")

add_text_box(s, Inches(0.8), Inches(1.2), Inches(10), Inches(0.8),
             "AIツール遍歴", font_size=36, color=INK, bold=True)
accent_bar(s, Inches(0.8), Inches(2.1))

# Timeline - left
timeline_data = [
    ("2022.11", "ChatGPT 3.5", "初めてのAI。飛びついた→飽きた"),
    ("2025 春", "GPT-4o 復帰", "ジブリ風で再燃。今度は毎日使う"),
    ("2025", "Cursor / Gemini", "コードエディタにAI統合。開発が変わった"),
    ("2026.01", "Claude との出会い", "「判断力」の時代だと気づく"),
    ("2026.03", "Claude Code 導入", "mdファイル1枚で会社を作れた"),
    ("2026.03", "AI幕府 始動", "エージェント12体。江戸幕府モデル"),
    ("2026.04", "現在", "放置で回る構造が完成しつつある"),
]

for i, (date, title, desc) in enumerate(timeline_data):
    y = Inches(2.5) + Inches(i * 0.65)
    # Date
    add_text_box(s, Inches(0.8), y, Inches(1.8), Inches(0.5),
                 date, font_size=13, color=ROUTE, bold=True,
                 font_name="JetBrains Mono")
    # Title
    add_text_box(s, Inches(2.8), y, Inches(3.2), Inches(0.5),
                 title, font_size=17, color=INK, bold=True)
    # Desc
    add_text_box(s, Inches(6.2), y, Inches(6.5), Inches(0.5),
                 desc, font_size=15, color=INK_SOFT)

# ====================================================================
# SLIDE 6: 転機 — Claude Code
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, MOSS_DEEP)
add_text_box(s, Inches(0.8), Inches(0.5), Inches(8), Inches(0.4),
             "CHAPTER 04 · TURNING POINT", font_size=11,
             color=ACCENT_GREEN, bold=True, font_name="JetBrains Mono")
add_slide_number(s, 7)
add_text_box(s, Inches(11.5), Inches(6.8), Inches(1.5), Inches(0.3),
             "~ 3 min", font_size=9, color=FADED,
             alignment=PP_ALIGN.RIGHT, font_name="JetBrains Mono")

add_text_box(s, Inches(0.8), Inches(1.2), Inches(10), Inches(1.2),
             "Claude Codeとの出会いが\nすべてを変えた", font_size=40,
             color=SNOW, bold=True)
accent_bar(s, Inches(0.8), Inches(2.8), Inches(2))

add_multiline(s, Inches(0.8), Inches(3.4), Inches(5.5), Inches(3.5), [
    "CLAUDE.md というファイル1枚に、",
    "会社の憲法を書いた。",
    "",
    "AIがそれを読んで、",
    "自律的に仕事を始めた。",
    "",
    "プログラミング初心者の僕が、",
    "AIの「上司」になった瞬間。",
], font_size=20, color=SNOW, spacing=Pt(4))

# Right side: Claude Code session style
card(s, Inches(7.2), Inches(3.0), Inches(5.3), Inches(3.5),
     bg_color=RGBColor(0x1B, 0x2A, 0x1F),
     border_color=RGBColor(0x4F, 0x6A, 0x48))
add_text_box(s, Inches(7.4), Inches(3.1), Inches(4.8), Inches(0.3),
             "$ claude", font_size=10, color=ACCENT_GREEN,
             font_name="JetBrains Mono")
add_multiline(s, Inches(7.4), Inches(3.5), Inches(4.9), Inches(2.8), [
    "> note\u8a18\u4e8b\u306e\u8349\u7a3f\u3092\u66f8\u3044\u3066",
    "",
    "\u2192 \u5c06\u8ecd\u304c\u6587\u658e\u3078\u30eb\u30fc\u30c6\u30a3\u30f3\u30b0...",
    "",
    "> \u3053\u306e\u30b3\u30fc\u30c9\u306e\u30d0\u30b0\u3092\u76f4\u3057\u3066",
    "",
    "\u2192 \u5c06\u8ecd\u304c\u9cf3\u5c71\u3078\u30eb\u30fc\u30c6\u30a3\u30f3\u30b0...",
    "",
    "> \u4eca\u9031\u306e\u4f7f\u7528\u91cf\u3092\u5831\u544a\u3057\u3066",
    "",
    "\u2192 \u5c06\u8ecd\u304c\u7b97\u4e4b\u52a9\u3078\u30eb\u30fc\u30c6\u30a3\u30f3\u30b0...",
], font_size=13, color=ACCENT_GREEN, font_name="JetBrains Mono", spacing=Pt(1))

# ====================================================================
# SLIDE 8: AI幕府の全体像 (4min)
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, PAPER)
add_section_label(s, "CHAPTER 04 · AI BAKUFU")
add_slide_number(s, 8)
add_text_box(s, Inches(11.5), Inches(6.8), Inches(1.5), Inches(0.3),
             "~ 4 min", font_size=9, color=INK_SOFT,
             alignment=PP_ALIGN.RIGHT, font_name="JetBrains Mono")

add_text_box(s, Inches(0.8), Inches(1.2), Inches(10), Inches(0.8),
             "AI幕府 — 12体のエージェント", font_size=36, color=INK, bold=True)
accent_bar(s, Inches(0.8), Inches(2.1))

# Top: Shogun (large card)
card(s, Inches(0.8), Inches(2.5), Inches(11.7), Inches(1.4),
     bg_color=MOSS_DEEP, border_color=MOSS_DEEP)
add_text_box(s, Inches(1.0), Inches(2.6), Inches(5), Inches(0.5),
             "\u5c06\u8ecd\uff08Claude Code\uff09", font_size=22, color=SNOW, bold=True)
add_text_box(s, Inches(1.0), Inches(3.15), Inches(10), Inches(0.5),
             "\u5e55\u5e9c\u7d71\u62ec\u3002\u30bf\u30b9\u30af\u3092\u53d7\u3051\u305f\u3089\u81ea\u5f8b\u5224\u65ad\u3067\u30c1\u30fc\u30e0\u7de8\u6210\u3057\u3001\u5fc5\u8981\u306a\u30a8\u30fc\u30b8\u30a7\u30f3\u30c8\u306b\u59d4\u4efb\u3059\u308b", font_size=14, color=RGBColor(0xCF, 0xD5, 0xB4))

# Middle: 5 key agents
key_agents = [
    ("\u6797 \u6587\u658e", "\u74e6\u7248\u65b9", "\u30b3\u30f3\u30c6\u30f3\u30c4\u30fb\u30de\u30fc\u30b1"),
    ("\u6797 \u9cf3\u5c71", "\u982d\u53d6", "\u958b\u767a\u30fb\u8abf\u67fb\u30fb\u6226\u7565"),
    ("\u672c\u591a \u667a\u623f", "\u8ecd\u5e2b", "\u91cd\u3044\u958b\u767a\u30fb\u8a2d\u8a08"),
    ("\u5411\u4e95 \u5f71\u7db1", "\u5fa1\u5ead\u756a", "\u60c5\u5831\u53ce\u96c6"),
    ("\u677e\u5e73 \u7b97\u4e4b\u52a9", "\u52d8\u5b9a\u5949\u884c", "\u4e88\u7b97\u30fb\u4f7f\u7528\u91cf\u7ba1\u7406"),
]

for i, (name, title, desc) in enumerate(key_agents):
    left = Inches(0.8 + i * 2.4)
    top = Inches(4.2)
    card(s, left, top, Inches(2.2), Inches(1.6))
    add_text_box(s, left + Inches(0.12), top + Inches(0.1),
                 Inches(2.0), Inches(0.35), name,
                 font_size=15, color=MOSS_DEEP, bold=True)
    add_text_box(s, left + Inches(0.12), top + Inches(0.5),
                 Inches(2.0), Inches(0.3), title,
                 font_size=11, color=ROUTE, font_name="JetBrains Mono")
    add_text_box(s, left + Inches(0.12), top + Inches(0.85),
                 Inches(2.0), Inches(0.6), desc,
                 font_size=12, color=INK_SOFT)

# Bottom: "+7 more" text
add_text_box(s, Inches(0.8), Inches(6.2), Inches(11), Inches(0.4),
             "+ \u4ed67\u4f53\uff08\u30bb\u30ad\u30e5\u30ea\u30c6\u30a3\u30fb\u54c1\u8cea\u76e3\u67fb\u30fbgit\u7ba1\u7406\u30fbKPI\u5206\u6790\u30fb\u30a4\u30f3\u30d5\u30e9\u30fb\u53cd\u8a3c\u30fb\u5be9\u8b70\uff09\u304c\u5f85\u6a5f\u4e2d",
             font_size=12, color=INK_SOFT, font_name="JetBrains Mono")

# ====================================================================
# SLIDE 8: 今の使い分け
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, PAPER)
add_section_label(s, "CHAPTER 04 · HOW I USE THEM")
add_slide_number(s, 9)
add_text_box(s, Inches(11.5), Inches(6.8), Inches(1.5), Inches(0.3),
             "~ 3 min", font_size=9, color=INK_SOFT,
             alignment=PP_ALIGN.RIGHT, font_name="JetBrains Mono")

add_text_box(s, Inches(0.8), Inches(1.2), Inches(10), Inches(0.8),
             "今の使い分け", font_size=36, color=INK, bold=True)
accent_bar(s, Inches(0.8), Inches(2.1))

tools = [
    ("Claude Code", "メイン。将軍として全タスクを統括\n12体のサブエージェントを自律運用", "常時稼働"),
    ("Claude Sonnet/Opus", "テキスト生成・推敲（Sonnet）\n重い思考・戦略検証（Opus）", "常時稼働"),
    ("Gemini Flash", "クイック検索・画像認識\n判断に迷った時のセカンドオピニオン", "随時"),
    ("Ollama（ローカル）", "qwen2.5:7b でA/B検証\nAPIコストゼロで実験し放題", "実験時"),
]

for i, (name, desc, status) in enumerate(tools):
    left = Inches(0.8 + (i % 2) * 6.2)
    top = Inches(2.5 + (i // 2) * 2.4)
    card(s, left, top, Inches(5.8), Inches(2.0))
    add_text_box(s, left + Inches(0.2), top + Inches(0.12),
                 Inches(4), Inches(0.4), name,
                 font_size=18, color=MOSS_DEEP, bold=True)
    add_text_box(s, left + Inches(4.2), top + Inches(0.12),
                 Inches(1.4), Inches(0.4), status,
                 font_size=10, color=ROUTE, font_name="JetBrains Mono",
                 alignment=PP_ALIGN.RIGHT)
    add_text_box(s, left + Inches(0.2), top + Inches(0.65),
                 Inches(5.2), Inches(1.2), desc,
                 font_size=14, color=INK_SOFT)

# ====================================================================
# SLIDE 9: 放置で回る構造
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, PAPER)
add_section_label(s, "CHAPTER 04 · AUTONOMOUS LOOP")
add_slide_number(s, 10)
add_text_box(s, Inches(11.5), Inches(6.8), Inches(1.5), Inches(0.3),
             "~ 3 min", font_size=9, color=INK_SOFT,
             alignment=PP_ALIGN.RIGHT, font_name="JetBrains Mono")

add_text_box(s, Inches(0.8), Inches(1.2), Inches(10), Inches(0.8),
             "放置で回る構造", font_size=36, color=INK, bold=True)
accent_bar(s, Inches(0.8), Inches(2.1))

add_multiline(s, Inches(0.8), Inches(2.6), Inches(11), Inches(4.5), [
    "Mac mini M4が24時間動いている。僕が寝ている間も。",
    "",
    "**2時間ごと** — 全市場スキャン＋自動トレード判定",
    "**15分ごと** — ポジション監視",
    "**30分ごと** — AI社員の死活チェック（heartbeat）",
    "**日2回** — 卒業条件の自動判定",
    "**月1回** — 資産配分の自動リバランス",
    "",
    "僕の仕事は「意思決定」と「承認ボタン」だけ。",
    "実行は全部AIがやる。",
], font_size=20, spacing=Pt(4))

# ====================================================================
# SLIDE 10: 失敗談
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, PAPER)
add_section_label(s, "CHAPTER 05 · FAILURES")
add_slide_number(s, 11)
add_text_box(s, Inches(11.5), Inches(6.8), Inches(1.5), Inches(0.3),
             "~ 3 min", font_size=9, color=INK_SOFT,
             alignment=PP_ALIGN.RIGHT, font_name="JetBrains Mono")

add_text_box(s, Inches(0.8), Inches(1.2), Inches(10), Inches(0.8),
             "失敗から学んだこと", font_size=36, color=INK, bold=True)
accent_bar(s, Inches(0.8), Inches(2.1))

# Failure 1
card(s, Inches(0.8), Inches(2.5), Inches(5.8), Inches(2.0))
add_text_box(s, Inches(1.0), Inches(2.6), Inches(5.4), Inches(0.4),
             "FAIL 01 — APIキー平文コミット", font_size=14,
             color=ROUTE, bold=True, font_name="JetBrains Mono")
add_multiline(s, Inches(1.0), Inches(3.1), Inches(5.4), Inches(1.2), [
    "セキュリティ報告書にAPIキーをそのまま書いた。",
    "git履歴に残るから、消しても手遅れ。",
    "→ 以降「先頭4文字 + ...」ルールを全社徹底",
], font_size=15, color=INK_SOFT)

# Failure 2
card(s, Inches(7.0), Inches(2.5), Inches(5.8), Inches(2.0))
add_text_box(s, Inches(7.2), Inches(2.6), Inches(5.4), Inches(0.4),
             "FAIL 02 — AIプロセス暴走", font_size=14,
             color=ROUTE, bold=True, font_name="JetBrains Mono")
add_multiline(s, Inches(7.2), Inches(3.1), Inches(5.4), Inches(1.2), [
    "Claude Max $200契約。エージェント12体を24時間稼働。",
    "大量に生成したけど、中身を見てなかった。",
    "→ AIは回すだけでは価値が出ない。回収して初めて価値になる",
], font_size=15, color=INK_SOFT)

# Failure 3
card(s, Inches(0.8), Inches(5.0), Inches(12), Inches(1.8))
add_text_box(s, Inches(1.0), Inches(5.1), Inches(11.5), Inches(0.4),
             "FAIL 03 — マイクラAIが全部忘れた", font_size=14,
             color=ROUTE, bold=True, font_name="JetBrains Mono")
add_multiline(s, Inches(1.0), Inches(5.6), Inches(11.5), Inches(1.0), [
    "AIに「マイクラで生きろ」と放り込んだ。拠点を4つ作った……のに全部忘れた。",
    "地図で見たら4つの点が散らばってた。「家に帰ろう」という感覚、AIにはない。",
    "→ AIは「手書き記録の外では記憶を活用できない」。だからCLAUDE.mdが必要。",
], font_size=15, color=INK_SOFT)

# ====================================================================
# SLIDE 11: AIで自分の仕事をどう変えるか
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, MOSS_DEEP)
add_text_box(s, Inches(0.8), Inches(0.5), Inches(8), Inches(0.4),
             "CHAPTER 06 · WHAT CAN WE BUILD?", font_size=11,
             color=ACCENT_GREEN, bold=True, font_name="JetBrains Mono")
add_slide_number(s, 12)
add_text_box(s, Inches(11.5), Inches(6.8), Inches(1.5), Inches(0.3),
             "~ 2 min", font_size=9, color=FADED,
             alignment=PP_ALIGN.RIGHT, font_name="JetBrains Mono")

add_text_box(s, Inches(0.8), Inches(1.3), Inches(10), Inches(1.2),
             "AIで自分の仕事を\nどう変えるか", font_size=40,
             color=SNOW, bold=True)
accent_bar(s, Inches(0.8), Inches(3.0), Inches(2))

add_multiline(s, Inches(0.8), Inches(3.6), Inches(11), Inches(3.5), [
    "便利なSaaSはたくさんある。でも社内導入はセキュリティ審査が厳しい。",
    "",
    "**だったら、自分で作ればいい。**",
    "",
    "AIがあれば3時間でMVPが作れる時代。",
    "コードレビュー要約、会議の構造化、リアルタイムマインドマップ——",
    "「欲しいけど導入できない」を、自分の手で内製する。",
    "",
    "今日のハッカソンで、1つ作ってみよう。",
], font_size=20, color=SNOW, spacing=Pt(6))

# ====================================================================
# SLIDE 12: これからの展望
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, PAPER)
add_section_label(s, "CHAPTER 07 · OUTLOOK")
add_slide_number(s, 13)
add_text_box(s, Inches(11.5), Inches(6.8), Inches(1.5), Inches(0.3),
             "~ 2 min", font_size=9, color=INK_SOFT,
             alignment=PP_ALIGN.RIGHT, font_name="JetBrains Mono")

add_text_box(s, Inches(0.8), Inches(1.2), Inches(10), Inches(0.8),
             "これからの使い方", font_size=36, color=INK, bold=True)
accent_bar(s, Inches(0.8), Inches(2.1))

add_multiline(s, Inches(0.8), Inches(2.6), Inches(11), Inches(4.5), [
    "3ヶ月やって分かったこと。",
    "",
    "**① 「生成」ではなく「回収」を設計する。**",
    "   AIは回すだけでは価値が出ない。残す問いと気づきを決める。",
    "",
    "**② 「何を作るか」より「どの問題を解くか」。**",
    "   自分が困ってないものは続かない。",
    "",
    "**③ SaaSを待たずに、自分で作る。**",
    "   セキュリティ審査を待つより、3時間で内製する方が早い。",
], font_size=20, spacing=Pt(3))

# ====================================================================
# SLIDE 13: 今日のゴール → ハッカソンへ
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, PAPER)
add_section_label(s, "TODAY'S GOAL")
add_slide_number(s, 14)
add_text_box(s, Inches(11.5), Inches(6.8), Inches(1.5), Inches(0.3),
             "~ 1 min", font_size=9, color=INK_SOFT,
             alignment=PP_ALIGN.RIGHT, font_name="JetBrains Mono")

add_text_box(s, Inches(0.8), Inches(1.2), Inches(10), Inches(0.8),
             "今日やること", font_size=36, color=INK, bold=True)
accent_bar(s, Inches(0.8), Inches(2.1))

add_multiline(s, Inches(0.8), Inches(2.6), Inches(11), Inches(4.5), [
    "**① 自分に合ったAIツールを1つ見つける**",
    "   全部使う必要はない。1つ「これ」を持って帰ればいい。",
    "",
    "**② チームで「動くもの」を作る**",
    "   完璧じゃなくていい。3時間で「動く」が最優先。",
    "",
    "**③ 「次に何をやるか」を決めて帰る**",
    "   今日の体験を、明日からの武器にする。",
    "",
    "",
    "では、午後のハッカソンで実際に作りましょう。",
], font_size=22, spacing=Pt(3))

# ====================================================================
# SLIDE 14: End
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(s, MOSS_DEEP)

add_text_box(s, Inches(0.8), Inches(2.2), Inches(10), Inches(2.5),
             "登る山を\n決めよう。", font_size=64,
             color=SNOW, bold=True)

accent_bar(s, Inches(0.8), Inches(5.0), Inches(2))

add_text_box(s, Inches(0.8), Inches(5.5), Inches(10), Inches(0.6),
             "質問・感想、なんでも。", font_size=20,
             color=RGBColor(0xCF, 0xD5, 0xB4))

add_text_box(s, Inches(0.8), Inches(6.3), Inches(10), Inches(0.5),
             "Deaf Engineers AI Workshop Vol.01 · 2026.05.31",
             font_size=14, color=FADED, font_name="JetBrains Mono")

# ====================================================================
# Save
# ====================================================================
out = os.path.join(os.path.dirname(__file__), "A1_presentation.pptx")
prs.save(out)
print(f"Saved: {out} ({TOTAL} slides)")
